"""
JARVIS Office Builder
Automated generation and manipulation of Excel, Word, and PowerPoint documents
"""

from typing import Optional, List, Dict, Any
from pathlib import Path


class OfficeBuilder:
    """Builder for Office documents."""
    
    def __init__(self, output_dir: str = "."):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
    
    def create_excel_report(
        self,
        filename: str,
        sheets: Dict[str, List[List[str]]],
        styles: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Create an Excel workbook with multiple sheets.
        
        Args:
            filename: Output file name (with .xlsx extension)
            sheets: Dict of sheet_name -> list of rows (each row is list of cells)
            styles: Optional styling configuration
        
        Returns:
            Path to created file
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet
            
            for sheet_name, data in sheets.items():
                ws = wb.create_sheet(sheet_name)
                for row_idx, row in enumerate(data, 1):
                    for col_idx, cell in enumerate(row, 1):
                        cell_obj = ws.cell(row=row_idx, column=col_idx)
                        cell_obj.value = cell
                        
                        # Apply header styling
                        if row_idx == 1 and styles:
                            cell_obj.font = Font(bold=True)
                            cell_obj.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                            cell_obj.font = Font(bold=True, color="FFFFFF")
            
            output_path = self.output_dir / filename
            wb.save(output_path)
            return str(output_path)
        except ImportError:
            return "Error: openpyxl not installed. Run: pip install openpyxl"
        except Exception as e:
            return f"Error creating Excel: {e}"
    
    def create_word_report(
        self,
        filename: str,
        title: str,
        sections: Dict[str, str]
    ) -> str:
        """
        Create a Word document with sections.
        
        Args:
            filename: Output file name (with .docx extension)
            title: Document title
            sections: Dict of section_name -> content_text
        
        Returns:
            Path to created file
        """
        try:
            from docx import Document
            from docx.shared import Pt, Inches, RGBColor
            
            doc = Document()
            doc.add_heading(title, 0)
            
            for section_name, content in sections.items():
                doc.add_heading(section_name, 1)
                doc.add_paragraph(content)
            
            output_path = self.output_dir / filename
            doc.save(output_path)
            return str(output_path)
        except ImportError:
            return "Error: python-docx not installed. Run: pip install python-docx"
        except Exception as e:
            return f"Error creating Word: {e}"
    
    def create_powerpoint_presentation(
        self,
        filename: str,
        title: str,
        slides: List[Dict[str, str]]
    ) -> str:
        """
        Create a PowerPoint presentation.
        
        Args:
            filename: Output file name (with .pptx extension)
            title: Presentation title
            slides: List of slide data (title, content)
        
        Returns:
            Path to created file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(54)
            title_frame.paragraphs[0].font.bold = True
            
            # Content slides
            for slide_data in slides:
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = slide_data.get("title", "")
                tf = content.text_frame
                tf.text = slide_data.get("content", "")
            
            output_path = self.output_dir / filename
            prs.save(output_path)
            return str(output_path)
        except ImportError:
            return "Error: python-pptx not installed. Run: pip install python-pptx"
        except Exception as e:
            return f"Error creating PowerPoint: {e}"


# Global office builder instance
_builder = None


def get_office_builder(output_dir: str = ".") -> OfficeBuilder:
    """Get or create an office builder instance."""
    global _builder
    if _builder is None:
        _builder = OfficeBuilder(output_dir)
    return _builder
