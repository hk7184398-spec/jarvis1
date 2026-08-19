"""
JARVIS PowerPoint Template Workflow
Template-based PowerPoint generation and automation
"""

from typing import Optional, Dict, List, Any
from pathlib import Path
from dataclasses import dataclass


@dataclass
class TemplateConfig:
    """Configuration for a PowerPoint template."""
    template_path: str
    output_dir: str
    theme_colors: Dict[str, str] = None
    fonts: Dict[str, str] = None
    
    def __post_init__(self):
        if self.theme_colors is None:
            self.theme_colors = {}
        if self.fonts is None:
            self.fonts = {}


class PPTTemplateWorkflow:
    """Manages PowerPoint template-based workflows."""
    
    def __init__(self):
        self.templates: Dict[str, TemplateConfig] = {}
        self.workflows: Dict[str, Dict[str, Any]] = {}
    
    def register_template(
        self,
        template_name: str,
        template_path: str,
        output_dir: str = "."
    ) -> bool:
        """Register a PowerPoint template."""
        if not Path(template_path).exists():
            return False
        
        config = TemplateConfig(
            template_path=template_path,
            output_dir=output_dir
        )
        self.templates[template_name] = config
        return True
    
    def create_from_template(
        self,
        template_name: str,
        output_filename: str,
        slide_data: List[Dict[str, Any]]
    ) -> Optional[str]:
        """
        Create a PowerPoint from a template.
        
        Args:
            template_name: Name of registered template
            output_filename: Output file name
            slide_data: List of slide configurations
        
        Returns:
            Path to created file or None if failed
        """
        template = self.templates.get(template_name)
        if not template:
            return None
        
        try:
            from pptx import Presentation
            
            # Load template
            prs = Presentation(template.template_path)
            
            # Clear existing slides (except first)
            while len(prs.slides) > 1:
                rId = prs.slides._sldIdLst[1].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[1]
            
            # Add slides from data
            for data in slide_data:
                slide_layout = prs.slide_layouts[data.get('layout', 1)]
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title if present
                if 'title' in data and slide.shapes.title:
                    slide.shapes.title.text = data['title']
                
                # Set content if present
                if 'content' in data and len(slide.placeholders) > 1:
                    content_shape = slide.placeholders[1]
                    if content_shape.has_text_frame:
                        content_shape.text = data['content']
            
            output_path = Path(template.output_dir) / output_filename
            prs.save(output_path)
            return str(output_path)
        except Exception as e:
            print(f"Error creating presentation: {e}")
            return None
    
    def apply_theme(
        self,
        template_name: str,
        colors: Optional[Dict[str, str]] = None,
        fonts: Optional[Dict[str, str]] = None
    ) -> bool:
        """Apply theme customizations to a template."""
        template = self.templates.get(template_name)
        if not template:
            return False
        
        if colors:
            template.theme_colors.update(colors)
        if fonts:
            template.fonts.update(fonts)
        
        return True
    
    def create_report_presentation(
        self,
        template_name: str,
        title: str,
        sections: Dict[str, str],
        output_filename: str
    ) -> Optional[str]:
        """Create a report presentation from sections."""
        slide_data = [
            {"title": title, "layout": 0}  # Title slide
        ]
        
        for section_title, content in sections.items():
            slide_data.append({
                "title": section_title,
                "content": content,
                "layout": 1
            })
        
        return self.create_from_template(template_name, output_filename, slide_data)
    
    def batch_create(
        self,
        template_name: str,
        batch_data: List[Dict[str, Any]]
    ) -> List[Optional[str]]:
        """
        Create multiple presentations from a template.
        
        Args:
            template_name: Template to use
            batch_data: List of {output_filename, slide_data} dicts
        
        Returns:
            List of output paths (or None for failed items)
        """
        results = []
        for item in batch_data:
            output = self.create_from_template(
                template_name,
                item['output_filename'],
                item['slide_data']
            )
            results.append(output)
        
        return results
    
    def save_workflow(self, workflow_name: str, config: Dict[str, Any]) -> None:
        """Save a workflow configuration."""
        self.workflows[workflow_name] = config
    
    def load_workflow(self, workflow_name: str) -> Optional[Dict[str, Any]]:
        """Load a workflow configuration."""
        return self.workflows.get(workflow_name)


# Global workflow instance
_workflow = None


def get_ppt_workflow() -> PPTTemplateWorkflow:
    """Get or create the global PowerPoint workflow."""
    global _workflow
    if _workflow is None:
        _workflow = PPTTemplateWorkflow()
    return _workflow
